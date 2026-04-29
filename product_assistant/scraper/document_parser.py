"""
Универсальный парсер документов: PDF, DOCX, PPTX.

Скачивает файл по URL (через прокси если задан) и извлекает текст.
Тип файла определяется по расширению в URL.
"""

import io
import os
import re
from urllib.parse import urlparse

import httpx
from loguru import logger

try:
    import pymupdf as fitz
except ImportError:
    try:
        import fitz
    except ImportError:
        fitz = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx"}


def find_document_links(page, page_url: str) -> list[dict]:
    """
    Находит ссылки на документы через Playwright JS-evaluate.
    Браузер резолвит абсолютный URL корректно (учитывает <base> тег).
    Возвращает список {'url': ..., 'title': ..., 'ext': ...}.
    """
    all_links = page.evaluate("""
        () => Array.from(document.querySelectorAll('a[href]'))
            .map(a => ({ url: a.href, title: a.textContent.trim() || a.href }))
    """)

    results = []
    seen = set()
    for link in all_links:
        url = link["url"]
        ext = _get_extension(url)
        if ext not in SUPPORTED_EXTENSIONS or url in seen:
            continue
        seen.add(url)
        results.append({"url": url, "title": link["title"] or url, "ext": ext})

    return results


def extract_document_text(doc_url: str, timeout: int = 30) -> str | None:
    """
    Скачивает документ и возвращает извлечённый текст.
    Поддерживает PDF, DOCX, PPTX.
    """
    ext = _get_extension(doc_url)
    if ext not in SUPPORTED_EXTENSIONS:
        logger.warning("Неподдерживаемый тип документа: {}", doc_url)
        return None

    content = _download(doc_url, timeout)
    if content is None:
        return None

    extractors = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".pptx": _extract_pptx,
    }
    return extractors[ext](content, doc_url)


# ------------------------------------------------------------------
# Внутренние функции
# ------------------------------------------------------------------

def _get_extension(url: str) -> str:
    path = urlparse(url).path.lower()
    for ext in SUPPORTED_EXTENSIONS:
        if path.endswith(ext):
            return ext
    return ""


def _download(url: str, timeout: int) -> bytes | None:
    proxy_url = os.environ.get("HTTPS_PROXY") or os.environ.get("HTTP_PROXY")
    try:
        with httpx.Client(
            proxy=proxy_url,
            timeout=timeout,
            headers={"User-Agent": "Mozilla/5.0 (compatible; ProductQABot/1.0)"},
            follow_redirects=True,
        ) as client:
            response = client.get(url)
            response.raise_for_status()
            return response.content
    except Exception as exc:
        logger.warning("Ошибка при скачивании {}: {}", url, exc)
        return None


def _clean(text: str) -> str:
    return re.sub(r'\n{3,}', '\n\n', text).strip()


def _extract_pdf(content: bytes, url: str) -> str | None:
    if fitz is None:
        logger.error("pymupdf не установлен: pip install pymupdf")
        return None
    try:
        doc = fitz.open(stream=io.BytesIO(content), filetype="pdf")
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        text = _clean(text)
        if len(text) < 50:
            logger.debug("PDF пустой или защищён: {}", url)
            return None
        return text
    except Exception as exc:
        logger.warning("Ошибка при разборе PDF {}: {}", url, exc)
        return None


def _extract_docx(content: bytes, url: str) -> str | None:
    if DocxDocument is None:
        logger.error("python-docx не установлен: pip install python-docx")
        return None
    try:
        doc = DocxDocument(io.BytesIO(content))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())
        # Таблицы
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    parts.append(row_text)
        text = _clean("\n".join(parts))
        if len(text) < 50:
            return None
        return text
    except Exception as exc:
        logger.warning("Ошибка при разборе DOCX {}: {}", url, exc)
        return None


def _extract_pptx(content: bytes, url: str) -> str | None:
    if Presentation is None:
        logger.error("python-pptx не установлен: pip install python-pptx")
        return None
    try:
        prs = Presentation(io.BytesIO(content))
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        slide_texts.append(line)
            if slide_texts:
                parts.append(f"[Слайд {slide_num}]\n" + "\n".join(slide_texts))
        text = _clean("\n\n".join(parts))
        if len(text) < 50:
            return None
        return text
    except Exception as exc:
        logger.warning("Ошибка при разборе PPTX {}: {}", url, exc)
        return None
